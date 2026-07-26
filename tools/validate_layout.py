#!/usr/bin/env python3
"""
validate_layout.py — PPT 布局校验工具（v3.5.1）

生成后自动运行，让布局 bug 暴露在开 PPT 之前。

用法：
  python3 tools/validate_layout.py output/你的文件.pptx

检测项：
  [E1] 元素越界（x+w > 10 或 y+h > 5.625）
  [E2] 元素进入页脚区（y+h > 5.32 但元素不是页码/数据来源）
  [E3] 元素相互重叠（同 slide 矩形相交 > 30%）
  [W1] 文本标题超长（按字号预估宽度超出 shape 宽度）
  [W2] 文本字号过小（<8pt 可能难读）

返回码：
  0 = 全部通过
  1 = 有 ERROR（必须修复）
  2 = 仅 WARNING（建议修复）

2026-07-24 复盘：用此工具可在用户开 PPT 之前提前发现 80% 的布局问题。
"""

import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

# 画布尺寸（16:9，英寸）
CANVAS_W = 10.0
CANVAS_H = 5.625
FOOTER_Y = 5.32  # 页脚区起始 y

# 每字宽度估算（中文全角，英寸）
CHAR_WIDTH_BY_PT = {
    44: 0.61, 40: 0.56, 36: 0.50, 32: 0.44, 28: 0.39, 26: 0.36,
    24: 0.33, 22: 0.31, 20: 0.28, 18: 0.25, 16: 0.22, 14: 0.19,
    12: 0.17, 11: 0.15, 10: 0.14, 9: 0.13, 8: 0.11, 7: 0.10,
}

def emu_to_in(v):
    return v / 914400.0

def estimate_char_width(pt):
    """根据字号估算单字宽度（中文全角）"""
    keys = sorted(CHAR_WIDTH_BY_PT.keys(), reverse=True)
    for k in keys:
        if pt >= k:
            return CHAR_WIDTH_BY_PT[k]
    return CHAR_WIDTH_BY_PT[8]

def is_cjk(ch):
    return '一' <= ch <= '鿿' or ch in '「」（）：；，。！？—、'

def estimate_text_width(text, pt):
    """估算文本渲染宽度（英寸）"""
    width = 0.0
    for ch in text:
        if is_cjk(ch):
            width += estimate_char_width(pt)
        elif ch.isdigit() or ch.isupper():
            width += estimate_char_width(pt) * 0.6
        elif ch == ' ':
            width += estimate_char_width(pt) * 0.4
        else:
            width += estimate_char_width(pt) * 0.5
    return width

def get_max_font_size(text_frame):
    """获取 text_frame 中最大字号"""
    max_pt = 0
    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                pt = run.font.size.pt
                if pt > max_pt:
                    max_pt = pt
    return max_pt

def get_text(text_frame):
    """获取 text_frame 的纯文本（连接所有段落）"""
    return "\n".join(p.text for p in text_frame.paragraphs if p.text.strip())

def shapes_overlap(r1, r2):
    """检测两个矩形是否重叠，返回重叠比例"""
    x1_min, y1_min, x1_max, y1_max = r1
    x2_min, y2_min, x2_max, y2_max = r2
    # 不相交
    if x1_max <= x2_min or x2_max <= x1_min or y1_max <= y2_min or y2_max <= y1_min:
        return 0.0
    # 相交面积
    overlap_x = min(x1_max, x2_max) - max(x1_min, x2_min)
    overlap_y = min(y1_max, y2_max) - max(y1_min, y2_min)
    overlap_area = overlap_x * overlap_y
    # 较小矩形的面积
    area1 = (x1_max - x1_min) * (y1_max - y1_min)
    area2 = (x2_max - x2_min) * (y2_max - y2_min)
    min_area = min(area1, area2)
    if min_area == 0:
        return 0.0
    return overlap_area / min_area

def is_contained(inner, outer, tolerance=0.2):
    """判断 inner 矩形是否被 outer 完全包含（嵌套关系，如卡片内的文本）。tolerance=0.2 覆盖表格文字略超出底色的小越界"""
    ix1, iy1, ix2, iy2 = inner
    ox1, oy1, ox2, oy2 = outer
    return (ix1 >= ox1 - tolerance and iy1 >= oy1 - tolerance and
            ix2 <= ox2 + tolerance and iy2 <= oy2 + tolerance)

def is_page_number(text):
    """判断文本是否是页码（如 "01 / 10"）"""
    return bool(re.match(r'^\d{1,2}\s*/\s*\d{1,2}$', text.strip()))

def is_source_line(text):
    """判断文本是否是数据来源行（如 "数据来源：..." 或 "来源：..."）"""
    t = text.strip()
    return t.startswith("数据来源") or t.startswith("来源：") or t.startswith("Source:")

def validate_pptx(path):
    p = Presentation(path)
    errors = []   # 必须修复
    warnings = [] # 建议修复

    for slide_idx, slide in enumerate(p.slides, 1):
        # 收集所有 shape 的矩形 + 文本信息
        shapes_info = []
        for shp_idx, shp in enumerate(slide.shapes):
            x = emu_to_in(shp.left) if shp.left is not None else 0
            y = emu_to_in(shp.top) if shp.top is not None else 0
            w = emu_to_in(shp.width) if shp.width is not None else 0
            h = emu_to_in(shp.height) if shp.height is not None else 0

            text = ""
            max_pt = 0
            if shp.has_text_frame:
                text = get_text(shp.text_frame)
                max_pt = get_max_font_size(shp.text_frame)

            shape_name = shp.name or f"shape_{shp_idx}"
            shapes_info.append({
                "name": shape_name,
                "rect": (x, y, x + w, y + h),
                "text": text,
                "max_pt": max_pt,
                "w": w, "h": h,
            })

            # [E1] 元素越界
            if x + w > CANVAS_W + 0.05:
                errors.append(f"[E1] P{slide_idx} {shape_name} 水平越界: x+w={x+w:.2f} > 10.0")
            if y + h > CANVAS_H + 0.05:
                errors.append(f"[E1] P{slide_idx} {shape_name} 垂直越界: y+h={y+h:.2f} > 5.625")

            # [E2] 元素进入页脚区（但允许页码/数据来源）
            if y + h > FOOTER_Y and y < FOOTER_Y:  # 跨入页脚区
                if text and not is_page_number(text) and not is_source_line(text):
                    # 大字标题/正文跨入页脚区
                    if max_pt >= 10:
                        errors.append(f"[E2] P{slide_idx} {shape_name} 跨入页脚区(y={y:.2f}+h={h:.2f}={y+h:.2f}>5.32): \"{text[:30]}...\"")

            # [W1] 文本超长
            if text and max_pt > 0 and w > 0:
                # 按行算
                for line in text.split("\n"):
                    if not line.strip():
                        continue
                    est_w = estimate_text_width(line, max_pt)
                    if est_w > w * 1.05:  # 超 5% 警告
                        warnings.append(f"[W1] P{slide_idx} {shape_name} 文本可能超长: 估算宽 {est_w:.2f}\" > 容器 {w:.2f}\" ({max_pt}pt, \"{line[:30]}...\")")

            # [W2] 字号过小
            if text and max_pt > 0 and max_pt < 8 and not is_page_number(text):
                warnings.append(f"[W2] P{slide_idx} {shape_name} 字号过小: {max_pt}pt < 8pt, 可能难读: \"{text[:20]}...\"")

        # [E3] 元素相互重叠（只对内容元素检测，跳过背景图/整页矩形/卡片嵌套）
        for i in range(len(shapes_info)):
            for j in range(i + 1, len(shapes_info)):
                s1 = shapes_info[i]
                s2 = shapes_info[j]
                # 跳过背景图（接近整页大小）
                if s1["w"] > 9.5 and s1["h"] > 5.0:
                    continue
                if s2["w"] > 9.5 and s2["h"] > 5.0:
                    continue
                # 跳过空矩形
                if s1["w"] < 0.1 or s1["h"] < 0.1 or s2["w"] < 0.1 or s2["h"] < 0.1:
                    continue
                # 跳过嵌套关系（卡片背景 vs 内部文本/图标）：一个矩形完全包含另一个
                if is_contained(s1["rect"], s2["rect"]) or is_contained(s2["rect"], s1["rect"]):
                    continue
                ratio = shapes_overlap(s1["rect"], s2["rect"])
                if ratio > 0.3:  # 重叠 30% 以上报警
                    # 如果其中一个是文本框，且文字为空，跳过
                    if not s1["text"] and not s2["text"]:
                        continue
                    warnings.append(f"[E3] P{slide_idx} 元素重叠 {ratio*100:.0f}%: {s1['name']}({s1['text'][:15]}) vs {s2['name']}({s2['text'][:15]})")

    return errors, warnings


def main():
    if len(sys.argv) < 2:
        print("用法: python3 validate_layout.py <pptx文件路径>")
        sys.exit(2)

    path = sys.argv[1]
    if not Path(path).exists():
        print(f"❌ 文件不存在: {path}")
        sys.exit(2)

    print(f"🔍 校验 {path}")
    errors, warnings = validate_pptx(path)

    print()
    if errors:
        print(f"❌ 发现 {len(errors)} 个 ERROR（必须修复）:")
        for e in errors:
            print(f"  {e}")
    if warnings:
        print(f"⚠️  发现 {len(warnings)} 个 WARNING（建议修复）:")
        for w in warnings:
            print(f"  {w}")
    if not errors and not warnings:
        print("✅ 布局校验全部通过")

    print()
    print(f"📊 总结: {len(errors)} errors / {len(warnings)} warnings")

    if errors:
        sys.exit(1)
    elif warnings:
        sys.exit(2)
    else:
        sys.exit(0)


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
style_extractor.py — PPT 风格提取器

从用户上传的 .pptx 模版中提取视觉风格，生成 styles/<name>/style.yaml，
并自动登记到 styles/registry.yaml，供 skill 按风格生成 PPT 时选用。

提取内容：
  - 主题色板（theme XML: dk1/lt1/dk2/lt2/accent1-6）
  - 字体（major/minor 拉丁 + 东亚字体）
  - 实际使用的文字色 / 填充色频次统计
  - 页面尺寸、图片密度、字号层级
  - 推导色板（背景深/浅、主色、强调色、文字色）

用法：
  # 单文件精细提取
  python3 tools/style_extractor.py <模版.pptx> --name corporate_blue \
      --keywords "央企,蓝色,商务"

  # 批量提取 uploads/ 目录全部 .pptx（--force 覆盖同名）
  python3 tools/style_extractor.py --batch [dir]

依赖：python-pptx, pyyaml（requirements.txt 已含）
"""

import argparse
import copy
import os
import re
import sys
import zipfile
from collections import Counter
from datetime import datetime, timezone, timedelta
from xml.etree import ElementTree as ET

SKILL_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
STYLES_DIR = os.path.join(SKILL_ROOT, "styles")
UPLOADS_DIR = os.path.join(SKILL_ROOT, "uploads")
REGISTRY_PATH = os.path.join(STYLES_DIR, "registry.yaml")

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
EMU_PER_INCH = 914400


# ---------------------------------------------------------------
# theme XML 解析
# ---------------------------------------------------------------
def _srgb_hex(elem):
    """从 a:srgbClr 或 a:sysClr 元素取 hex（大写，不带#）。"""
    if elem is None:
        return None
    srgb = elem.find(f"{{{A_NS}}}srgbClr")
    if srgb is not None:
        return srgb.get("val", "").upper() or None
    sysclr = elem.find(f"{{{A_NS}}}sysClr")
    if sysclr is not None:
        return (sysclr.get("lastClr") or "").upper() or None
    return None


def extract_theme(pptx_path):
    """解析 ppt/theme/theme1.xml，返回 theme_colors 与 fonts。"""
    theme_colors, fonts = {}, {}
    with zipfile.ZipFile(pptx_path) as z:
        theme_names = [n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml", n)]
        if not theme_names:
            return theme_colors, fonts
        root = ET.fromstring(z.read(theme_names[0]))

    scheme = root.find(f".//{{{A_NS}}}clrScheme")
    if scheme is not None:
        for child in scheme:
            tag = child.tag.split("}")[-1]  # dk1/lt1/accent1...
            hexval = _srgb_hex(child)
            if hexval:
                theme_colors[tag] = hexval

    font_scheme = root.find(f".//{{{A_NS}}}fontScheme")
    if font_scheme is not None:
        for grp in ("majorFont", "minorFont"):
            node = font_scheme.find(f"{{{A_NS}}}{grp}")
            if node is None:
                continue
            latin = node.find(f"{{{A_NS}}}latin")
            ea = node.find(f"{{{A_NS}}}ea")
            fonts[f"{grp}_latin"] = latin.get("typeface") if latin is not None else None
            fonts[f"{grp}_ea"] = ea.get("typeface") if ea is not None else None
    return theme_colors, fonts


# ---------------------------------------------------------------
# 幻灯片内容统计
# ---------------------------------------------------------------
def collect_usage_stats(pptx_path):
    """用 python-pptx 统计实际使用的颜色/字体/图片密度。"""
    from pptx import Presentation
    from pptx.enum.dml import MSO_FILL_TYPE
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    text_colors, fill_colors, font_faces, title_sizes, body_sizes = Counter(), Counter(), Counter(), Counter(), Counter()
    total_images = 0
    pages_with_images = 0

    for slide in prs.slides:
        slide_images = 0
        for shape in slide.shapes:
            # 图片统计
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_images += 1
            # 形状填充色
            try:
                if shape.fill.type == MSO_FILL_TYPE.SOLID:
                    rgb = shape.fill.fore_color.rgb
                    if rgb is not None:
                        fill_colors[str(rgb).upper()] += 1
            except Exception:
                pass
            # 文本
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        font_faces[run.font.name] += 1
                    try:
                        if run.font.color and run.font.color.rgb is not None:
                            text_colors[str(run.font.color.rgb).upper()] += 1
                    except Exception:
                        pass
                    if run.font.size:
                        pt = round(run.font.size.pt)
                        if pt >= 20:
                            title_sizes[pt] += 1
                        elif pt >= 8:
                            body_sizes[pt] += 1
        total_images += slide_images
        if slide_images:
            pages_with_images += 1

    n = max(len(prs.slides), 1)
    return {
        "slide_count": len(prs.slides),
        "width_in": round(prs.slide_width / EMU_PER_INCH, 3),
        "height_in": round(prs.slide_height / EMU_PER_INCH, 3),
        "text_colors": text_colors,
        "fill_colors": fill_colors,
        "font_faces": font_faces,
        "title_sizes": title_sizes,
        "body_sizes": body_sizes,
        "total_images": total_images,
        "avg_images_per_slide": round(total_images / n, 1),
        "pages_with_images": pages_with_images,
    }


# ---------------------------------------------------------------
# 推导色板
# ---------------------------------------------------------------
def _hex_brightness(hexv):
    r, g, b = int(hexv[0:2], 16), int(hexv[2:4], 16), int(hexv[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b


def _is_grayish(hexv, tol=18):
    r, g, b = int(hexv[0:2], 16), int(hexv[2:4], 16), int(hexv[4:6], 16)
    return max(r, g, b) - min(r, g, b) <= tol


def derive_palette(theme_colors, stats):
    """从主题色+使用频次推导角色化色板。"""
    pal = {}

    dk1 = theme_colors.get("dk1")
    lt1 = theme_colors.get("lt1")
    pal["background_dark"] = dk1 or "1E293B"
    pal["background_light"] = lt1 or "FFFFFF"
    pal["card_dark"] = theme_colors.get("dk2", pal["background_dark"])
    pal["card_light"] = theme_colors.get("lt2", "FFFFFF")

    # 主色：accent1；若 accent1 偏灰则退而取使用频次最高的非灰填充色
    primary = theme_colors.get("accent1")
    if primary and _is_grayish(primary):
        for hexv, _ in stats["fill_colors"].most_common():
            if not _is_grayish(hexv):
                primary = hexv
                break
    pal["primary"] = primary or "0D9488"
    pal["primary_light"] = theme_colors.get("accent2", pal["primary"])
    pal["primary_bright"] = theme_colors.get("accent3", pal["primary_light"])
    pal["accent"] = theme_colors.get("accent4", "F59E0B")
    pal["accent2"] = theme_colors.get("accent5", "3B82F6")
    pal["success"] = "10B981"
    pal["warning"] = "EF4444"

    # 文字色：使用频次最高的深/浅色文字
    dark_texts = [h for h, _ in stats["text_colors"].most_common() if _hex_brightness(h) < 110]
    light_texts = [h for h, _ in stats["text_colors"].most_common() if _hex_brightness(h) > 200]
    pal["text_main"] = dark_texts[0] if dark_texts else "334155"
    pal["text_sub"] = dark_texts[1] if len(dark_texts) > 1 else "64748B"
    pal["text_on_dark"] = light_texts[0] if light_texts else "FFFFFF"
    pal["text_on_dark_sub"] = "94A3B8"
    return pal


# ---------------------------------------------------------------
# style.yaml 生成 + registry 登记
# ---------------------------------------------------------------
def _load_template_style():
    tpl_path = os.path.join(STYLES_DIR, "_template", "style.yaml")
    import yaml
    with open(tpl_path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def build_style_yaml(name, pptx_path, keywords, theme_colors, theme_fonts, stats):
    style = _load_template_style()

    style["name"] = name
    style["display_name"] = name
    style["source"] = f"extracted:{os.path.basename(pptx_path)}"
    style["extracted_at"] = datetime.now(timezone(timedelta(hours=8))).isoformat(timespec="seconds")
    style["keywords"] = keywords

    style["slide"]["width_in"] = stats["width_in"]
    style["slide"]["height_in"] = stats["height_in"]
    ratio = stats["width_in"] / stats["height_in"] if stats["height_in"] else 16 / 9
    style["slide"]["aspect"] = "16:9" if abs(ratio - 16 / 9) < 0.05 else ("4:3" if abs(ratio - 4 / 3) < 0.05 else f"{ratio:.2f}")

    style["palette"] = derive_palette(theme_colors, stats)

    latin = theme_fonts.get("majorFont_latin") or (stats["font_faces"].most_common(1)[0][0] if stats["font_faces"] else "Arial")
    cjk = theme_fonts.get("majorFont_ea") or "Noto Sans SC"
    style["fonts"]["latin"] = latin
    style["fonts"]["cjk"] = cjk
    if stats["title_sizes"]:
        style["fonts"]["title_max_pt"] = max(stats["title_sizes"])
        style["fonts"]["title_min_pt"] = min(stats["title_sizes"])
    if stats["body_sizes"]:
        style["fonts"]["body_pt"] = stats["body_sizes"].most_common(1)[0][0]

    style["extraction"] = {
        "theme_colors": theme_colors,
        "top_text_colors": [{"hex": h, "count": c} for h, c in stats["text_colors"].most_common(8)],
        "top_fill_colors": [{"hex": h, "count": c} for h, c in stats["fill_colors"].most_common(8)],
        "font_faces_seen": [f for f, _ in stats["font_faces"].most_common(8)],
        "avg_images_per_slide": stats["avg_images_per_slide"],
        "slide_count": stats["slide_count"],
    }
    return style


def register_style(name, keywords, description=""):
    """把新风格登记进 styles/registry.yaml（幂等：同名则更新）。"""
    import yaml
    with open(REGISTRY_PATH, "r", encoding="utf-8") as f:
        reg = yaml.safe_load(f)

    entry = {
        "name": name,
        "display_name": name,
        "path": f"{name}/style.yaml",
        "keywords": keywords,
        "description": description or f"提取自用户上传模版",
        "builtin": False,
    }
    styles = [s for s in reg.get("styles", []) if s.get("name") != name]
    styles.append(entry)
    reg["styles"] = styles
    reg["updated"] = datetime.now(timezone(timedelta(hours=8))).date().isoformat()

    with open(REGISTRY_PATH, "w", encoding="utf-8") as f:
        yaml.safe_dump(reg, f, allow_unicode=True, sort_keys=False, default_flow_style=False)


def _kw_hit(kw, text):
    """关键词命中判定（逐级放宽）：
    1. 完整子串命中
    2. 归一化命中：忽略大小写与 -_/空格
    3. 拉丁系关键词按 token 拆分，命中 ≥2 个即算（应对 "Ocean-Blue-Report" vs "ocean blue"）
    4. 含 CJK 字符的长关键词，任意 3 字连续片段命中（应对 "红色党建汇报模版"）
    """
    if kw in text:
        return True
    norm = lambda s: re.sub(r"[\s\-_]+", "", s).lower()
    nkw, ntext = norm(kw), norm(text)
    if nkw and nkw in ntext:
        return True
    has_cjk = any("一" <= ch <= "鿿" for ch in kw)
    if not has_cjk:
        tokens = [t for t in re.split(r"[^a-zA-Z0-9]+", kw) if t]
        if len(tokens) >= 2:
            hits = sum(1 for t in tokens if t.lower() in ntext)
            if hits >= 2:
                return True
        return False
    need = min(len(kw), 3)
    return any(kw[i:i + need] in text for i in range(len(kw) - need + 1))


def select_style(user_input, registry_path=REGISTRY_PATH):
    """按用户输入匹配风格（供其他模块调用）。返回 (style_name, style_dict)。"""
    import yaml
    with open(registry_path, "r", encoding="utf-8") as f:
        reg = yaml.safe_load(f)

    best, best_hits = None, 0
    for s in reg.get("styles", []):
        hits = sum(1 for kw in s.get("keywords", []) if kw and _kw_hit(kw, user_input))
        if hits > best_hits:
            best, best_hits = s, hits
        elif hits == best_hits and hits > 0 and best is not None:
            # 平手时优先用户提取的风格（比内置更贴近用户语境）
            if best.get("builtin") and not s.get("builtin"):
                best = s
    if best is None:
        fallback = reg.get("selection", {}).get("fallback")
        best = next((s for s in reg["styles"] if s["name"] == fallback), reg["styles"][0])

    style_path = os.path.join(STYLES_DIR, best["path"])
    with open(style_path, "r", encoding="utf-8") as f:
        return best["name"], yaml.safe_load(f)


def extract_one(pptx_path, name, keywords, description=""):
    """单文件提取全链路：解析→统计→生成style.yaml→登记registry。返回 style.yaml 路径。"""
    theme_colors, theme_fonts = extract_theme(pptx_path)
    stats = collect_usage_stats(pptx_path)
    style = build_style_yaml(name, pptx_path, keywords, theme_colors, theme_fonts, stats)
    out_dir = os.path.join(STYLES_DIR, name)
    os.makedirs(out_dir, exist_ok=True)
    import yaml
    out_path = os.path.join(out_dir, "style.yaml")
    with open(out_path, "w", encoding="utf-8") as f:
        yaml.safe_dump(style, f, allow_unicode=True, sort_keys=False, default_flow_style=False)
    register_style(name, keywords, description)
    return out_path, style


def _slugify_name(stem, fallback):
    """文件名 → snake_case 风格ID；无法转写时用 fallback。"""
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", stem).strip("_").lower()
    slug = re.sub(r"_+", "_", slug)
    if not slug or not slug[0].isalpha():
        return fallback
    return slug


def batch_extract(uploads_dir=UPLOADS_DIR, force=False):
    """批量扫描目录中的 .pptx，逐个提取风格入库。

    - 风格ID：文件名转 snake_case；中文等无法转写时命名 style_001...
    - keywords：自动纳入原始文件名（不含扩展名）便于中文名匹配
    - 已存在同名风格且未指定 force 时跳过
    """
    import yaml
    if not os.path.isdir(uploads_dir):
        print(f"[ERROR] 目录不存在: {uploads_dir}", file=sys.stderr)
        sys.exit(1)

    pptx_files = sorted(
        f for f in os.listdir(uploads_dir)
        if f.lower().endswith(".pptx") and not f.startswith("~$")
    )
    if not pptx_files:
        print(f"目录中无 .pptx 文件: {uploads_dir}")
        return

    with open(REGISTRY_PATH, "r", encoding="utf-8") as f:
        reg_styles = yaml.safe_load(f).get("styles", [])
    existing = {s["name"] for s in reg_styles}
    # 中文文件名走 style_NNN 序列命名，无法靠名字判重；
    # 用 keywords 中的原始文件名（stem）判重，保证 --batch 幂等
    existing_stems = {kw for s in reg_styles for kw in s.get("keywords", [])}

    ok, skipped, failed = [], [], []
    auto_seq = 1
    for fname in pptx_files:
        stem = os.path.splitext(fname)[0]

        if stem in existing_stems and not force:
            prev = next((s["name"] for s in reg_styles if stem in s.get("keywords", [])), "?")
            skipped.append((fname, prev))
            continue

        name = _slugify_name(stem, None)
        if name is None:
            while f"style_{auto_seq:03d}" in existing:
                auto_seq += 1
            name = f"style_{auto_seq:03d}"
            auto_seq += 1

        if name in existing and not force:
            skipped.append((fname, name))
            continue

        path = os.path.join(uploads_dir, fname)
        # keywords 保留原文件名分词（中文名直接整串纳入，匹配时按子串命中）
        keywords = [stem]
        try:
            out_path, style = extract_one(path, name, keywords,
                                          description=f"批量提取自 uploads/{fname}")
            existing.add(name)
            ok.append((fname, name, style["palette"]["primary"]))
        except Exception as e:
            failed.append((fname, str(e)))

    print(f"\n批量提取完成: 成功 {len(ok)}, 跳过 {len(skipped)}, 失败 {len(failed)}")
    for fname, name, primary in ok:
        print(f"  ✓ {fname} → {name} (主色 {primary})")
    for fname, name in skipped:
        print(f"  - {fname} → {name} 已存在，跳过（--force 可覆盖）")
    for fname, err in failed:
        print(f"  ✗ {fname}: {err}")


# ---------------------------------------------------------------
# CLI
# ---------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser(description="从用户上传的PPT模版提取风格样例")
    ap.add_argument("pptx", nargs="?", help="用户上传的 .pptx 模版路径（单文件模式）")
    ap.add_argument("--name", help="风格ID（snake_case，将作为 styles/ 子目录名）")
    ap.add_argument("--keywords", default="", help="逗号分隔的触发关键词，如 '央企,蓝色,商务'")
    ap.add_argument("--description", default="", help="风格描述（写入registry）")
    ap.add_argument("--batch", nargs="?", const=UPLOADS_DIR, default=None,
                    help="批量模式：扫描指定目录（默认 uploads/）全部 .pptx 并逐个提取")
    ap.add_argument("--force", action="store_true", help="批量模式下覆盖已存在的同名风格")
    args = ap.parse_args()

    if args.batch is not None:
        batch_extract(args.batch, force=args.force)
        return

    if not args.pptx or not args.name:
        ap.error("单文件模式需要 <pptx> 与 --name；批量模式用 --batch [dir]")

    if not os.path.isfile(args.pptx):
        print(f"[ERROR] 文件不存在: {args.pptx}", file=sys.stderr)
        sys.exit(1)
    if not re.match(r"^[a-z][a-z0-9_]*$", args.name):
        print(f"[ERROR] --name 必须是 snake_case: {args.name}", file=sys.stderr)
        sys.exit(1)

    keywords = [k.strip() for k in args.keywords.split(",") if k.strip()]

    print(f"[1/2] 提取风格: {args.pptx}")
    out_path, style = extract_one(args.pptx, args.name, keywords, args.description)
    print(f"[2/2] 已登记 registry.yaml")

    print(f"\n✓ 风格 '{args.name}' 已入库: {out_path}")
    print(f"  主色: {style['palette']['primary']}  深底: {style['palette']['background_dark']}  浅底: {style['palette']['background_light']}")
    print(f"  字体: {style['fonts']['latin']} / {style['fonts']['cjk']}")


if __name__ == "__main__":
    main()

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, Any
import os
import re


# ================ 参考PPT风格定义 ================
# 基于《职业规划完整指南》PPT逆向提取的视觉规范
STYLE = {
    # 页面尺寸：16:9 宽屏
    "slide_width": 13.33,
    "slide_height": 7.50,

    # 配色系统（Tailwind Emerald + Orange accent）
    "dark_bg": "#1E293B",       # 封面深色背景
    "primary": "#047857",       # 主色-深绿
    "primary_mid": "#059669",   # 主色-中绿
    "primary_light": "#10B981", # 主色-亮绿
    "accent": "#6EE7B7",        # 强调-浅绿
    "orange": "#F97316",        # 橙色强调
    "orange_dark": "#EA580C",   # 深橙
    "orange_tint_bg": "#FFF7ED", # 橙色浅底
    "green_tint_bg": "#E6F9F0", # 绿色浅底
    "card_bg": "#FFFFFF",       # 卡片白色
    "text_dark": "#1F2937",     # 正文深色
    "text_mid": "#374151",      # 次级文字
    "text_gray": "#6B7280",     # 辅助文字
    "text_light": "#94A3B8",    # 浅色文字
    "border_light": "#E5E7EB",  # 浅边框
    "page_bg": "#FFFFFF",       # 页面背景

    # 字体
    "font": "Noto Sans SC",

    # 字号层级
    "cover_title": 42,       # 封面大标题
    "cover_subtitle": 24,    # 封面副标题
    "cover_caption": 18,     # 封面说明
    "section_label": 28,     # PART标签
    "section_title": 44,     # 章节标题
    "section_subtitle": 22,  # 章节副标题
    "page_title": 32,        # 页面标题
    "item_header": 18,       # 条目标题
    "item_desc": 14,         # 条目描述
    "card_label": 16,        # 卡片标签
    "card_desc": 13,         # 卡片描述
    "quote": 14,             # 引述
    "badge_num": 20,         # 编号数字
    "qa_title": 48,          # Q&A标题

    # 间距
    "margin_left": 1.11,
    "margin_right": 1.11,
    "content_width": 11.11,
}


class PPTXGenerator:
    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.vi = config.get("vi_spec", {})
        self.s = STYLE  # 风格快捷引用

    def _rgb(self, hex_color: str) -> RGBColor:
        """十六进制转RGBColor"""
        return RGBColor.from_string(hex_color.lstrip("#"))

    def generate_pptx(self, page_content: Dict[str, Any], layout_spec: Dict[str, Any], materials: Dict[str, Any]) -> str:
        """生成可编辑PPT文件，参照职业规划PPT风格"""
        s = self.s
        prs = Presentation()
        prs.slide_width = Inches(s["slide_width"])
        prs.slide_height = Inches(s["slide_height"])

        first_key = next(iter(page_content), "presentation")
        cover_title = page_content.get(first_key, {}).get("title", "presentation")
        output_dir = "output"
        os.makedirs(output_dir, exist_ok=True)
        safe_title = cover_title.replace("\uff1a", "-").replace("/", "-").replace("\\", "-")
        output_path = os.path.join(output_dir, f"{safe_title}.pptx")

        for page_id, content in page_content.items():
            layout_type = content.get("layout_type", "bullets")
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            if layout_type == "cover":
                self._render_cover(slide, content)
            elif layout_type == "toc":
                self._render_toc(slide, content)
            elif layout_type == "steps":
                self._render_steps(slide, content)
            elif layout_type == "matrix":
                self._render_matrix(slide, content)
            elif layout_type == "comparison":
                self._render_comparison(slide, content)
            elif layout_type == "closing":
                self._render_closing(slide, content)
            else:
                self._render_bullets(slide, content, page_id, layout_spec, materials)

            if content.get("speaker_notes"):
                self._add_speaker_notes(slide, content["speaker_notes"])

        prs.save(output_path)
        print(f"✅ PPT文件已保存：{output_path}")
        return output_path

    # ================ 封面页 ================
    def _render_cover(self, slide, content):
        s = self.s
        ml = s["margin_left"]
        cw = s["content_width"]

        # 深色全屏背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(s["slide_width"]), Inches(s["slide_height"]))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(s["dark_bg"])
        bg.line.fill.background()

        # 主标题（居中偏上）
        title = content.get("title", "")
        self._add_text(slide, title,
                       x=1.5, y=1.8, w=10.33, h=1.5,
                       font_size=s["cover_title"], bold=True, color=s["accent"],
                       alignment=PP_ALIGN.CENTER)

        # 副标题
        subtitle = content.get("subtitle", "")
        if subtitle:
            self._add_text(slide, subtitle,
                           x=1.5, y=3.3, w=10.33, h=0.8,
                           font_size=s["cover_subtitle"], bold=False, color=s["page_bg"],
                           alignment=PP_ALIGN.CENTER)

        # 装饰横线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(4.3), Inches(3.33), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = self._rgb(s["accent"])
        line.line.fill.background()

        # 底部三个关键指标卡片
        cards = [
            ("效率 +50%", "运营效率提升"),
            ("成本 -40%", "人力成本优化"),
            ("AUM 翻倍", "资产管理规模"),
        ]
        card_w = 3.06
        card_h = 1.25
        card_gap = 0.5
        total_w = len(cards) * card_w + (len(cards) - 1) * card_gap
        start_x = (s["slide_width"] - total_w) / 2
        card_y = 5.2

        for i, (label, desc) in enumerate(cards):
            cx = start_x + i * (card_w + card_gap)
            # 卡片白底
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(card_y), Inches(card_w), Inches(card_h))
            card.fill.solid()
            card.fill.fore_color.rgb = self._rgb(s["card_bg"])
            card.line.fill.background()

            # 标签
            self._add_text(slide, label,
                           x=cx + 0.15, y=card_y + 0.15, w=card_w - 0.3, h=0.5,
                           font_size=16, bold=True, color=s["primary_light"],
                           alignment=PP_ALIGN.CENTER)
            # 描述
            self._add_text(slide, desc,
                           x=cx + 0.15, y=card_y + 0.65, w=card_w - 0.3, h=0.4,
                           font_size=12, bold=False, color=s["text_gray"],
                           alignment=PP_ALIGN.CENTER)

    # ================ 目录页 ================
    def _render_toc(self, slide, content):
        s = self.s
        ml = s["margin_left"]

        # 白色页面背景
        self._add_page_bg(slide)

        # 标题
        self._add_text(slide, content.get("title", ""),
                       x=ml, y=0.6, w=s["content_width"], h=0.8,
                       font_size=s["page_title"], bold=True, color=s["primary"],
                       alignment=PP_ALIGN.LEFT)

        # 分隔线
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ml), Inches(1.45), Inches(s["content_width"]), Inches(0.03))
        sep.fill.solid()
        sep.fill.fore_color.rgb = self._rgb(s["accent"])
        sep.line.fill.background()

        # 目录项（卡片式）
        body = content.get("body", "")
        if body:
            items = [line.strip() for line in body.split("\n") if line.strip()]
            start_y = 1.8
            card_h = 0.85
            gap = 0.2

            for i, item in enumerate(items):
                item_text = item.lstrip("0123456789. ").strip()
                cy = start_y + i * (card_h + gap)

                # 卡片白底
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ml), Inches(cy), Inches(s["content_width"]), Inches(card_h))
                card.fill.solid()
                card.fill.fore_color.rgb = self._rgb(s["card_bg"])
                card.line.color.rgb = self._rgb(s["border_light"])
                card.line.width = Pt(1)

                # 编号圆圈
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ml + 0.2), Inches(cy + 0.17), Inches(0.5), Inches(0.5))
                circle.fill.solid()
                circle.fill.fore_color.rgb = self._rgb(s["primary"])
                circle.line.fill.background()

                ctf = circle.text_frame
                ctf.word_wrap = False
                cp = ctf.paragraphs[0]
                cp.text = str(i + 1)
                cp.font.name = s["font"]
                cp.font.size = Pt(s["badge_num"])
                cp.font.bold = True
                cp.font.color.rgb = self._rgb(s["card_bg"])
                cp.alignment = PP_ALIGN.CENTER
                ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 目录文字
                self._add_text(slide, item_text,
                               x=ml + 0.9, y=cy + 0.15, w=s["content_width"] - 1.2, h=0.55,
                               font_size=20, bold=True, color=s["text_dark"])

    # ================ 要点页 ================
    def _render_bullets(self, slide, content, page_id, layout_spec, materials):
        s = self.s
        ml = s["margin_left"]

        # 白色页面背景
        self._add_page_bg(slide)

        # 标题
        self._add_text(slide, content.get("title", ""),
                       x=ml, y=0.6, w=s["content_width"], h=0.8,
                       font_size=s["page_title"], bold=True, color=s["primary"],
                       alignment=PP_ALIGN.LEFT)

        # 分隔线
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ml), Inches(1.45), Inches(s["content_width"]), Inches(0.03))
        sep.fill.solid()
        sep.fill.fore_color.rgb = self._rgb(s["accent"])
        sep.line.fill.background()

        # 判断是否有图表
        has_chart = page_id in materials.get("charts", {})
        has_image = page_id in materials.get("images", {})

        # 要点列表（卡片式）
        body = content.get("body", "")
        if body:
            items = [line.strip() for line in body.split("\n") if line.strip()]

            if has_chart or has_image:
                # 左文右图布局
                bullet_w = 6.5
                chart_x = 8.0
                chart_w = 4.5
            else:
                # 全宽布局
                bullet_w = s["content_width"]

            start_y = 1.7
            card_h = 1.0
            gap = 0.15

            for i, item in enumerate(items):
                item_text = item.lstrip("\u2022\u25b8 ").strip()
                if not item_text:
                    continue

                cy = start_y + i * (card_h + gap)

                # 卡片白底
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ml), Inches(cy), Inches(bullet_w), Inches(card_h))
                card.fill.solid()
                card.fill.fore_color.rgb = self._rgb(s["card_bg"])
                card.line.color.rgb = self._rgb(s["border_light"])
                card.line.width = Pt(1)

                # 编号徽章（橙色圆角矩形）
                badge_w = 0.5
                badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ml + 0.15), Inches(cy + 0.2), Inches(badge_w), Inches(0.5))
                badge.fill.solid()
                badge.fill.fore_color.rgb = self._rgb(s["orange"])
                badge.line.fill.background()

                btf = badge.text_frame
                btf.word_wrap = False
                bp = btf.paragraphs[0]
                bp.text = str(i + 1)
                bp.font.name = s["font"]
                bp.font.size = Pt(16)
                bp.font.bold = True
                bp.font.color.rgb = self._rgb(s["card_bg"])
                bp.alignment = PP_ALIGN.CENTER
                btf.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 要点文字
                self._add_text(slide, item_text,
                               x=ml + 0.8, y=cy + 0.15, w=bullet_w - 1.0, h=0.7,
                               font_size=s["item_desc"], bold=False, color=s["text_dark"])

        # 图表区域
        if has_chart:
            chart_path = materials["charts"][page_id]
            if os.path.exists(chart_path):
                try:
                    slide.shapes.add_picture(chart_path, Inches(chart_x), Inches(1.7), width=Inches(chart_w))
                except Exception as e:
                    print(f"  ⚠️ 图表嵌入失败：{e}")

        if has_image and not has_chart:
            image_path = materials["images"][page_id]
            if os.path.exists(image_path):
                try:
                    slide.shapes.add_picture(image_path, Inches(chart_x), Inches(1.7), width=Inches(chart_w))
                except Exception as e:
                    print(f"  ⚠️ 图片嵌入失败：{e}")

    # ================ 步骤页 ================
    def _render_steps(self, slide, content):
        s = self.s
        ml = s["margin_left"]

        # 白色页面背景
        self._add_page_bg(slide)

        # 标题
        self._add_text(slide, content.get("title", ""),
                       x=ml, y=0.6, w=s["content_width"], h=0.8,
                       font_size=s["page_title"], bold=True, color=s["primary"],
                       alignment=PP_ALIGN.LEFT)

        # 分隔线
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ml), Inches(1.45), Inches(s["content_width"]), Inches(0.03))
        sep.fill.solid()
        sep.fill.fore_color.rgb = self._rgb(s["accent"])
        sep.line.fill.background()

        # 步骤卡片
        body = content.get("body", "")
        if body:
            items = [line.strip() for line in body.split("\n") if line.strip()]
            start_y = 1.7
            card_h = 1.15
            gap = 0.2

            step_colors = [s["primary"], s["primary_mid"], s["primary_light"], s["accent"]]
            step_labels = ["F", "A", "S", "T"]

            for i, item in enumerate(items):
                item_text = item.lstrip("\u25b8\u2022 ").strip()
                if not item_text:
                    continue

                cy = start_y + i * (card_h + gap)
                step_color = step_colors[i % len(step_colors)]

                # 步骤标签（左侧彩色条）
                tag_w = 0.7
                tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ml), Inches(cy), Inches(tag_w), Inches(card_h))
                tag.fill.solid()
                tag.fill.fore_color.rgb = self._rgb(step_color)
                tag.line.fill.background()

                ttf = tag.text_frame
                ttf.word_wrap = False
                tp = ttf.paragraphs[0]
                tp.text = step_labels[i] if i < 4 else str(i + 1)
                tp.font.name = s["font"]
                tp.font.size = Pt(28)
                tp.font.bold = True
                tp.font.color.rgb = self._rgb(s["card_bg"])
                tp.alignment = PP_ALIGN.CENTER
                ttf.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 内容卡片
                card_w = s["content_width"] - tag_w - 0.1
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ml + tag_w + 0.1), Inches(cy), Inches(card_w), Inches(card_h))
                card.fill.solid()
                card.fill.fore_color.rgb = self._rgb(s["card_bg"])
                card.line.color.rgb = self._rgb(step_color)
                card.line.width = Pt(2)

                ctf = card.text_frame
                ctf.word_wrap = True
                cp = ctf.paragraphs[0]
                cp.text = item_text
                cp.font.name = s["font"]
                cp.font.size = Pt(s["item_desc"])
                cp.font.color.rgb = self._rgb(s["text_dark"])
                ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
                ctf.margin_left = Inches(0.2)
                ctf.margin_right = Inches(0.2)

                # 连接箭头
                if i < len(items) - 1:
                    arrow_y = cy + card_h
                    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(ml + 0.2), Inches(arrow_y), Inches(0.3), Inches(gap))
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = self._rgb(s["primary_light"])
                    arrow.line.fill.background()

        # 底部总结条
        self._add_bottom_bar(slide, "24个月转型路线：Foundation → Agent → Switch → Transcend")

    # ================ 2x2矩阵页 ================
    def _render_matrix(self, slide, content):
        """2×2四象限布局（SWOT/BCG/Agent矩阵）

        content 结构：
          title: 页面标题
          quadrants: [{"label": "S 优势", "items": ["..."], "color": "#10B981"(可选)}, ...] 最多4个
          summary: 底部总结条文字（可选）
          body: 无 quadrants 时的回退，每行格式 "标签|要点1;要点2"
        """
        s = self.s
        ml = s["margin_left"]

        self._add_page_bg(slide)

        # 标题 + 分隔线（与其他内容页一致）
        self._add_text(slide, content.get("title", ""),
                       x=ml, y=0.6, w=s["content_width"], h=0.8,
                       font_size=s["page_title"], bold=True, color=s["primary"],
                       alignment=PP_ALIGN.LEFT)
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ml), Inches(1.45), Inches(s["content_width"]), Inches(0.03))
        sep.fill.solid()
        sep.fill.fore_color.rgb = self._rgb(s["accent"])
        sep.line.fill.background()

        quadrants = content.get("quadrants")
        if not quadrants:
            quadrants = []
            for line in content.get("body", "").split("\n"):
                line = line.strip().lstrip("\u2022\u25b8 ").strip()
                if not line:
                    continue
                if "|" in line:
                    label, items_str = line.split("|", 1)
                    items = [i.strip() for i in re.split(r"[;；]", items_str) if i.strip()]
                else:
                    label, items = line, []
                quadrants.append({"label": label, "items": items})
        quadrants = quadrants[:4]
        if not quadrants:
            return

        q_colors = [s["primary"], s["orange"], s["primary_mid"], s["primary_light"]]
        area_top, area_bottom = 1.7, 6.1
        gap = 0.25
        qw = (s["content_width"] - gap) / 2
        qh = (area_bottom - area_top - gap) / 2

        for i, q in enumerate(quadrants):
            row, col = divmod(i, 2)
            qx = ml + col * (qw + gap)
            qy = area_top + row * (qh + gap)
            q_color = q.get("color") or q_colors[i % len(q_colors)]

            # 象限卡片：白底 + 彩色描边
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(qx), Inches(qy), Inches(qw), Inches(qh))
            card.fill.solid()
            card.fill.fore_color.rgb = self._rgb(s["card_bg"])
            card.line.color.rgb = self._rgb(q_color)
            card.line.width = Pt(2)

            # 象限标签
            self._add_text(slide, q.get("label", f"Q{i+1}"),
                           x=qx + 0.25, y=qy + 0.12, w=qw - 0.5, h=0.4,
                           font_size=s["card_label"], bold=True, color=q_color)

            # 象限要点
            items = q.get("items") or []
            max_items = max(int((qh - 0.7) / 0.32), 1)
            for j, item in enumerate(items[:max_items]):
                self._add_text(slide, "\u2022 " + str(item),
                               x=qx + 0.3, y=qy + 0.6 + j * 0.32, w=qw - 0.6, h=0.3,
                               font_size=s["card_desc"], bold=False, color=s["text_dark"])

        if content.get("summary"):
            self._add_bottom_bar(slide, content["summary"])

    # ================ 对比分析页 ================
    def _render_comparison(self, slide, content):
        """左右对比布局（传统 vs 新范式 / 现状 vs 目标）

        content 结构：
          title: 页面标题
          comparison: {"left": {"title": "传统模式", "items": [...]},
                       "right": {"title": "新范式", "items": [...]}}
          summary: 底部总结条文字（可选）
          body: 无 comparison 时的回退，每行格式 "左侧要点|右侧要点"
        """
        s = self.s
        ml = s["margin_left"]

        self._add_page_bg(slide)

        self._add_text(slide, content.get("title", ""),
                       x=ml, y=0.6, w=s["content_width"], h=0.8,
                       font_size=s["page_title"], bold=True, color=s["primary"],
                       alignment=PP_ALIGN.LEFT)
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ml), Inches(1.45), Inches(s["content_width"]), Inches(0.03))
        sep.fill.solid()
        sep.fill.fore_color.rgb = self._rgb(s["accent"])
        sep.line.fill.background()

        comp = content.get("comparison") or {}
        left = comp.get("left") or {}
        right = comp.get("right") or {}
        if not left and not right:
            left_items, right_items = [], []
            for line in content.get("body", "").split("\n"):
                line = line.strip().lstrip("\u2022\u25b8 ").strip()
                if not line:
                    continue
                if "|" in line:
                    l, r = line.split("|", 1)
                    left_items.append(l.strip())
                    right_items.append(r.strip())
                else:
                    left_items.append(line)
            left = {"title": "现状 / 传统", "items": left_items}
            right = {"title": "目标 / 新范式", "items": right_items}

        gap = 0.4
        col_w = (s["content_width"] - gap) / 2
        area_top, card_h = 1.7, 4.4
        columns = [
            {"data": left,  "x": ml,                    "head_color": s["text_gray"], "bg": s["card_bg"],      "border": s["border_light"]},
            {"data": right, "x": ml + col_w + gap,      "head_color": s["primary"],   "bg": s["green_tint_bg"], "border": s["primary"]},
        ]

        for col in columns:
            data, cx = col["data"], col["x"]

            # 列卡片
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(area_top), Inches(col_w), Inches(card_h))
            card.fill.solid()
            card.fill.fore_color.rgb = self._rgb(col["bg"])
            card.line.color.rgb = self._rgb(col["border"])
            card.line.width = Pt(1.5)

            # 列标题（彩色 pill）
            pill_w = min(len(str(data.get("title", ""))) * 0.22 + 0.6, col_w - 0.5)
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx + 0.25), Inches(area_top + 0.2), Inches(pill_w), Inches(0.5))
            pill.fill.solid()
            pill.fill.fore_color.rgb = self._rgb(col["head_color"])
            pill.line.fill.background()
            ptf = pill.text_frame
            ptf.word_wrap = False
            pp = ptf.paragraphs[0]
            pp.text = str(data.get("title", ""))
            pp.font.name = s["font"]
            pp.font.size = Pt(s["card_label"])
            pp.font.bold = True
            pp.font.color.rgb = self._rgb("#FFFFFF")
            pp.alignment = PP_ALIGN.CENTER
            ptf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 列要点
            items = data.get("items") or []
            max_items = max(int((card_h - 1.0) / 0.42), 1)
            for j, item in enumerate(items[:max_items]):
                self._add_text(slide, "\u2022 " + str(item),
                               x=cx + 0.3, y=area_top + 0.9 + j * 0.42, w=col_w - 0.6, h=0.4,
                               font_size=s["item_desc"], bold=False, color=s["text_dark"])

        # 中间 VS 标识
        vs_x = ml + col_w + gap / 2 - 0.35
        vs = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(vs_x), Inches(area_top + card_h / 2 - 0.35), Inches(0.7), Inches(0.7))
        vs.fill.solid()
        vs.fill.fore_color.rgb = self._rgb(s["orange"])
        vs.line.fill.background()
        vtf = vs.text_frame
        vtf.word_wrap = False
        vp = vtf.paragraphs[0]
        vp.text = "VS"
        vp.font.name = s["font"]
        vp.font.size = Pt(16)
        vp.font.bold = True
        vp.font.color.rgb = self._rgb("#FFFFFF")
        vp.alignment = PP_ALIGN.CENTER
        vtf.vertical_anchor = MSO_ANCHOR.MIDDLE

        if content.get("summary"):
            self._add_bottom_bar(slide, content["summary"])

    # ================ 结尾页 ================
    def _render_closing(self, slide, content):
        s = self.s

        # 深色全屏背景（同封面）
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(s["slide_width"]), Inches(s["slide_height"]))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(s["dark_bg"])
        bg.line.fill.background()

        # 主标题居中
        title = content.get("title", "")
        self._add_text(slide, title,
                       x=1.5, y=1.8, w=10.33, h=1.5,
                       font_size=s["qa_title"], bold=True, color=s["accent"],
                       alignment=PP_ALIGN.CENTER)

        # 副文本
        body = content.get("body", "")
        if body:
            lines = [l.strip() for l in body.split("\n") if l.strip()]
            first_line = lines[0] if lines else ""
            self._add_text(slide, first_line,
                           x=2.0, y=3.3, w=9.33, h=0.8,
                           font_size=s["cover_caption"], bold=False, color=s["text_light"],
                           alignment=PP_ALIGN.CENTER)

        # Q&A
        self._add_text(slide, "Q & A",
                       x=3.0, y=4.5, w=7.33, h=1.0,
                       font_size=s["qa_title"], bold=True, color=s["primary_light"],
                       alignment=PP_ALIGN.CENTER)

        # 底部装饰线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(5.8), Inches(4.33), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = self._rgb(s["accent"])
        line.line.fill.background()

        # 感谢文字
        self._add_text(slide, "感谢聆听与关注",
                       x=3.0, y=6.0, w=7.33, h=0.6,
                       font_size=s["cover_caption"], bold=False, color=s["text_light"],
                       alignment=PP_ALIGN.CENTER)

    # ================ 公共组件 ================
    def _add_page_bg(self, slide):
        """白色页面背景"""
        s = self.s
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(s["slide_width"]), Inches(s["slide_height"]))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(s["page_bg"])
        bg.line.fill.background()

    def _add_bottom_bar(self, slide, text):
        """底部绿色总结条"""
        s = self.s
        bar_y = 6.3
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(s["margin_left"]), Inches(bar_y), Inches(s["content_width"]), Inches(0.7))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb(s["accent"])
        bar.line.fill.background()

        self._add_text(slide, text,
                       x=s["margin_left"] + 0.2, y=bar_y + 0.1, w=s["content_width"] - 0.4, h=0.5,
                       font_size=14, bold=True, color=s["primary"],
                       alignment=PP_ALIGN.LEFT)

    def _add_text(self, slide, text, x, y, w, h, font_size=18, bold=False, color="#1F2937", alignment=PP_ALIGN.LEFT):
        """通用文本渲染"""
        s = self.s
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = s["font"]
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = self._rgb(color)
        p.alignment = alignment

    def _add_speaker_notes(self, slide, notes_text):
        """写入演讲备注"""
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

    # ================ 兼容旧接口 ================
    def _add_text_zone(self, slide, zone: Dict[str, Any], text: str):
        x, y, w, h = zone["x"], zone["y"], zone["w"], zone["h"]
        font_name = zone.get("font", STYLE["font"])
        font_size = zone.get("font_size", 18)
        color_hex = zone.get("color", "#1F2937")
        self._add_text(slide, text, x, y, w, h, font_size, False, color_hex)

    def _add_subtitle(self, slide, title_zone: Dict[str, Any], sub_y: float, subtitle: str):
        x = title_zone.get("x", 1.11)
        w = title_zone.get("w", 11.11)
        self._add_text(slide, subtitle, x, sub_y, w, 0.6, 16, False, "#6B7280")

    def _add_code_zone(self, slide, zone: Dict[str, Any], code: str):
        x, y, w, h = zone["x"], zone["y"], zone["w"], zone["h"]
        bg_color = zone.get("bg_color", "#1E293B")
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(bg_color.lstrip("#"))
        shape.line.fill.background()
        self._add_text(slide, code, x + 0.1, y + 0.1, w - 0.2, h - 0.2, 14, False, "#E2E8F0")

    def _add_image_zone(self, slide, zone: Dict[str, Any], image_path: str):
        x, y, w, h = zone["x"], zone["y"], zone["w"], zone["h"]
        if os.path.exists(image_path):
            try:
                slide.shapes.add_picture(image_path, Inches(x), Inches(y), width=Inches(w))
            except Exception as e:
                print(f"  ⚠️ 图片嵌入失败({image_path})：{e}，跳过")
        else:
            print(f"  ⚠️ 图片文件不存在：{image_path}，跳过嵌入")

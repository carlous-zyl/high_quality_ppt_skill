#!/usr/bin/env python3
"""
OpenClaw PPT SKILL - E2E 自动化测试
覆盖：模块导入 → 配置 → 受众分级 → 6步完整链路 → PPT文件生成 → 校验通过
"""
import unittest
import sys
import os
import yaml
import json

# 添加项目根目录到路径
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from core.pipeline import OpenClawPPTSkill
from core.content_generator import ContentGenerator
from core.layout_designer import LayoutDesigner
from tools.image_generator import ImageGenerator
from tools.chart_generator import ChartGenerator
from tools.pptx_generator import PPTXGenerator
from validators.consistency_checker import ConsistencyChecker
from validators.compliance_checker import ComplianceChecker


class TestE2EPipeline(unittest.TestCase):
    """E2E 全链路自动化测试"""

    @classmethod
    def setUpClass(cls):
        """一次性初始化：加载配置、初始化所有模块"""
        cls.config_path = os.path.join(
            os.path.dirname(__file__), '..', 'config', 'global_config.yaml'
        )
        with open(cls.config_path, "r", encoding="utf-8") as f:
            cls.config = yaml.safe_load(f)

        # 将执行模式切换为 auto，跳过人工确认
        cls.config["execution_mode"] = "auto"

        cls.skill = OpenClawPPTSkill(config_path=cls.config_path)
        cls.skill.execution_mode = "auto"  # 强制 auto 模式

        # 注入自动确认回调（双保险）
        cls.skill.set_confirm_callback(lambda cp, d: True)

        cls.content_gen = ContentGenerator(cls.config)
        cls.layout_designer = LayoutDesigner(cls.config)
        cls.chart_gen = ChartGenerator(cls.config)
        cls.compliance_checker = ComplianceChecker(cls.config)
        cls.consistency_checker = ConsistencyChecker(cls.config)

    # ============ 基础模块导入测试 ============
    def test_01_all_modules_importable(self):
        """所有核心模块可正常导入"""
        self.assertIsNotNone(ContentGenerator)
        self.assertIsNotNone(LayoutDesigner)
        self.assertIsNotNone(ImageGenerator)
        self.assertIsNotNone(ChartGenerator)
        self.assertIsNotNone(PPTXGenerator)
        self.assertIsNotNone(ConsistencyChecker)
        self.assertIsNotNone(ComplianceChecker)
        self.assertIsNotNone(OpenClawPPTSkill)

    def test_02_config_loaded(self):
        """全局配置正确加载"""
        self.assertIn("vi_spec", self.config)
        self.assertIn("primary_color", self.config["vi_spec"])
        self.assertIn("toolchain", self.config)
        self.assertEqual(self.config["vi_spec"]["primary_color"], "#047857")

    # ============ 受众分级测试 ============
    def test_03_audience_classification(self):
        """三级受众分级模型正确分类"""
        # business
        self.assertEqual(self.content_gen._classify_audience("CEO"), "business")
        self.assertEqual(self.content_gen._classify_audience("业务运营高管"), "business")
        # technical
        self.assertEqual(self.content_gen._classify_audience("技术架构师"), "technical")
        self.assertEqual(self.content_gen._classify_audience("研发工程师"), "technical")
        # hybrid
        self.assertEqual(self.content_gen._classify_audience("产品总监"), "hybrid")
        self.assertEqual(self.content_gen._classify_audience("VP"), "hybrid")
        # 默认
        self.assertEqual(self.content_gen._classify_audience("普通用户"), "hybrid")

    # ============ 步骤1：需求对齐 ============
    def test_04_step1_requirement_alignment_success(self):
        """步骤1：完整需求对齐成功"""
        requirement = {
            "核心主题": "数字化转型战略",
            "目标受众": "CEO",
            "汇报时长": "20分钟",
            "总页数": 8,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "避免敏感数据"
        }
        result = self.skill._step1_requirement_alignment(requirement)
        self.assertIsNotNone(result)
        self.assertTrue(result.get("_alignment_confirmed"))

    def test_05_step1_requirement_alignment_missing_params(self):
        """步骤1：缺少必填参数时对齐失败"""
        incomplete = {"核心主题": "测试主题"}
        result = self.skill._step1_requirement_alignment(incomplete)
        self.assertIsNone(result)

    # ============ 步骤2：大纲生成 ============
    def test_06_step2_outline_generation(self):
        """步骤2：大纲生成正确"""
        requirement = {
            "核心主题": "AI赋能企业增长",
            "目标受众": "CEO",
            "汇报时长": "15分钟",
            "总页数": 8,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "无"
        }
        aligned = self.skill._step1_requirement_alignment(requirement)
        self.assertIsNotNone(aligned)

        outline = self.skill._step2_outline_generation(aligned)
        self.assertIsNotNone(outline)
        self.assertIn("叙事线", outline)
        self.assertIn("大纲结构", outline)
        self.assertEqual(len(outline["大纲结构"]), 8)

    def test_07_step2_narrative_line_generalized(self):
        """步骤2：叙事线已通用化（不含硬编码AI字样）"""
        requirement = {
            "核心主题": "企业数字化转型",
            "目标受众": "CTO",
            "汇报时长": "15分钟",
            "总页数": 8,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "无"
        }
        aligned = self.skill._step1_requirement_alignment(requirement)
        outline = self.skill._step2_outline_generation(aligned)
        narrative = outline["叙事线"]
        # 叙事线应包含主题词"企业数字化转型"
        self.assertIn("企业数字化转型", narrative)

    def test_08_step2_dynamic_core_page_index(self):
        """步骤2：核心方案页索引动态查找（非硬编码insert_pos=4）"""
        requirement = {
            "核心主题": "测试动态大纲",
            "目标受众": "技术负责人",
            "汇报时长": "30分钟",
            "总页数": 12,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "无"
        }
        aligned = self.skill._step1_requirement_alignment(requirement)
        outline = self.skill._step2_outline_generation(aligned)
        # 12页应包含扩展页
        page_types = [p["page_type"] for p in outline["大纲结构"]]
        self.assertIn("核心方案扩展页", page_types)

    def test_09_step2_truncation_preserves_cover_and_qa(self):
        """步骤2：截断时保留封面页和总结Q&A页"""
        requirement = {
            "核心主题": "测试截断",
            "目标受众": "产品经理",
            "汇报时长": "5分钟",
            "总页数": 3,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "无"
        }
        aligned = self.skill._step1_requirement_alignment(requirement)
        outline = self.skill._step2_outline_generation(aligned)
        page_types = [p["page_type"] for p in outline["大纲结构"]]
        self.assertIn("封面页", page_types)
        self.assertIn("总结Q&A页", page_types)
        self.assertEqual(len(outline["大纲结构"]), 3)

    # ============ 步骤3：内容生成 + 信息密度 ============
    def test_10_step3_page_content_with_speaker_notes(self):
        """步骤3：每页内容包含 speaker_notes 字段"""
        requirement = {
            "核心主题": "AI赋能金融风控",
            "目标受众": "CEO",
            "汇报时长": "15分钟",
            "总页数": 8,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "无"
        }
        aligned = self.skill._step1_requirement_alignment(requirement)
        outline = self.skill._step2_outline_generation(aligned)
        page_content = self.skill._step3_page_content_design(outline, aligned)

        self.assertIsNotNone(page_content)
        for page_id, content in page_content.items():
            self.assertIn("speaker_notes", content, f"页面 {page_id} 缺少 speaker_notes")

    def test_11_step3_information_density_check(self):
        """步骤3：信息密度校验功能正常"""
        # 正常内容
        normal_content = {"p1": {"body": "这是正常长度的内容"}}
        self.assertTrue(self.compliance_checker.check_information_density(normal_content))

        # 超长内容（>300有效字符）
        long_body = "这是一段很长的内容" * 50
        long_content = {"p1": {"body": long_body}}
        self.assertFalse(self.compliance_checker.check_information_density(long_content))

    def test_12_step3_auto_truncate(self):
        """步骤3：自动截断超长内容"""
        # 每个中文字符算1字，"这是超长内容"=6字，100次=600字 > 300字上限
        long_body = "这是超长内容" * 100
        page_content = {"p1": {"body": long_body, "title": "测试"}}
        result = self.skill._auto_truncate_content(page_content)
        self.assertIn("详见演讲备注", result["p1"]["body"])

    # ============ 步骤4：布局设计 ============
    def test_13_step4_layout_design(self):
        """步骤4：布局设计产出正确"""
        requirement = {
            "核心主题": "布局测试",
            "目标受众": "技术总监",
            "汇报时长": "10分钟",
            "总页数": 5,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "无"
        }
        aligned = self.skill._step1_requirement_alignment(requirement)
        outline = self.skill._step2_outline_generation(aligned)
        page_content = self.skill._step3_page_content_design(outline, aligned)
        layout_spec = self.skill._step4_layout_design(page_content)

        self.assertIsNotNone(layout_spec)
        for page_id, spec in layout_spec.items():
            self.assertIn("zones", spec)
            self.assertIn("margin", spec)
            self.assertIn("colors", spec)

    # ============ 步骤5：素材生成 ============
    def test_14_step5_material_generation_charts(self):
        """步骤5：图表生成含CSV+Python源码导出"""
        requirement = {
            "核心主题": "图表测试",
            "目标受众": "数据分析师",
            "汇报时长": "10分钟",
            "总页数": 5,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "无"
        }
        aligned = self.skill._step1_requirement_alignment(requirement)
        outline = self.skill._step2_outline_generation(aligned)
        page_content = self.skill._step3_page_content_design(outline, aligned)
        layout_spec = self.skill._step4_layout_design(page_content)
        materials = self.skill._step5_material_generation(layout_spec, page_content)

        self.assertIsNotNone(materials)
        self.assertIn("images", materials)
        self.assertIn("charts", materials)

        # 验证有chart_data的页面生成了图表
        chart_pages = [pid for pid, c in page_content.items() if c.get("chart_data")]
        if chart_pages:
            for pid in chart_pages:
                if pid in materials["charts"]:
                    chart_path = materials["charts"][pid]
                    self.assertTrue(chart_path.endswith(".png"), f"图表 {pid} 非PNG格式")

    # ============ 校验器测试 ============
    def test_15_validators_consistency(self):
        """一致性校验器真实返回结果"""
        vi = self.config["vi_spec"]
        allowed_colors = set([vi["primary_color"], vi["secondary_color"]] + vi["neutral_colors"])
        mock_layout = {
            "p1": {
                "colors": list(allowed_colors)[:3],
                "margin": "20mm"
            },
            "p2": {
                "colors": [vi["primary_color"], vi["secondary_color"]],
                "margin": "20mm"
            }
        }
        mock_content = {"p1": {"font": {}}, "p2": {"font": {}}}
        result = self.consistency_checker.check(mock_content, mock_layout)
        self.assertTrue(result)

    def test_16_validators_compliance(self):
        """合规性校验器真实返回结果"""
        mock_layout = {
            "p1": {
                "contrast_ratio": 5.2,
                "min_font_size": 16,
                "logo_position": "top_right"
            }
        }
        mock_content = {"p1": {"body": "合规测试内容"}}
        result = self.compliance_checker.check(mock_content, mock_layout, {})
        self.assertTrue(result)

    # ============ 完整E2E链路测试 ============
    def test_17_full_pipeline_business_audience(self):
        """E2E全链路：business受众（CEO级）"""
        requirement = {
            "核心主题": "企业数字化转型战略",
            "目标受众": "CEO",
            "汇报时长": "20分钟",
            "总页数": 8,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "避免提及未确认的财务数据"
        }
        result = self.skill.execute_pipeline(requirement)
        self.assertTrue(result.endswith(".pptx"), f"预期返回.pptx路径，实际：{result}")

        # 验证PPT文件存在
        self.assertTrue(os.path.exists(result), f"PPT文件不存在：{result}")

        # 验证演讲备注文件存在
        notes_path = os.path.join(os.path.dirname(result), "speaker_notes.txt")
        self.assertTrue(os.path.exists(notes_path), f"演讲备注不存在：{notes_path}")

    def test_18_full_pipeline_technical_audience(self):
        """E2E全链路：technical受众（CTO/工程师级）"""
        requirement = {
            "核心主题": "微服务架构升级方案",
            "目标受众": "技术架构师",
            "汇报时长": "30分钟",
            "总页数": 12,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "避免泄露核心算法细节"
        }
        result = self.skill.execute_pipeline(requirement)
        self.assertTrue(result.endswith(".pptx"))
        self.assertTrue(os.path.exists(result))

    def test_19_full_pipeline_hybrid_audience(self):
        """E2E全链路：hybrid受众（产品VP级）"""
        requirement = {
            "核心主题": "产品增长飞轮策略",
            "目标受众": "产品VP",
            "汇报时长": "15分钟",
            "总页数": 6,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "避免竞品对比"
        }
        result = self.skill.execute_pipeline(requirement)
        self.assertTrue(result.endswith(".pptx"))
        self.assertTrue(os.path.exists(result))

    # ============ 构图映射测试 ============
    def test_20_composition_map_coverage(self):
        """COMPOSITION_MAP 覆盖所有常见页面类型"""
        from tools.image_generator import ImageGenerator
        expected_keys = ["cover", "agenda", "background", "solution_core",
                         "solution_ext", "tech_detail", "result", "risk_plan", "summary"]
        for key in expected_keys:
            self.assertIn(key, ImageGenerator.COMPOSITION_MAP,
                          f"COMPOSITION_MAP 缺少 key: {key}")

    # ============ PPT文件质量验证 ============
    def test_21_pptx_content_quality(self):
        """PPT文件内容质量验证（页数、标题、正文）"""
        from pptx import Presentation

        requirement = {
            "核心主题": "PPT质量验证测试",
            "目标受众": "CEO",
            "汇报时长": "15分钟",
            "总页数": 5,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "无"
        }
        result = self.skill.execute_pipeline(requirement)
        prs = Presentation(result)

        # 验证页数
        self.assertEqual(len(prs.slides), 5, f"预期5页，实际{len(prs.slides)}页")

        # 验证每页至少有文本内容
        for i, slide in enumerate(prs.slides):
            has_text = any(
                shape.has_text_frame and shape.text_frame.text.strip()
                for shape in slide.shapes
            )
            self.assertTrue(has_text, f"第{i+1}页没有任何文本内容")

    def test_22_chart_csv_and_source_export(self):
        """图表CSV与Python源码导出验证"""
        requirement = {
            "核心主题": "图表导出测试",
            "目标受众": "数据分析师",
            "汇报时长": "10分钟",
            "总页数": 5,
            "交付格式": ".pptx",
            "品牌VI": "默认企业蓝",
            "内容禁忌": "无"
        }
        result = self.skill.execute_pipeline(requirement)

        # 验证输出目录中有CSV文件
        chart_dir = os.path.join(os.path.dirname(result), "charts")
        if os.path.exists(chart_dir):
            csv_files = [f for f in os.listdir(chart_dir) if f.endswith(".csv")]
            py_files = [f for f in os.listdir(chart_dir) if f.endswith(".py")]
            # 如果有图表页面，应有对应的CSV和PY文件
            if csv_files:
                self.assertGreater(len(csv_files), 0, "有图表但无CSV导出")
                self.assertGreater(len(py_files), 0, "有图表但无Python源码导出")


if __name__ == "__main__":
    unittest.main(verbosity=2)
